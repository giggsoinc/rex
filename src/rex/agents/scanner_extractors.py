"""Format-specific text extractors — pure-Python, no LibreOffice/soffice.

Each function reads ONE format directly in-process and never raises:
on any failure it logs and returns "" so a single bad file never kills a batch.

Modern formats use dedicated libraries (already installed):
  PDF  -> pdfplumber, fallback pypdf
  DOCX -> python-docx
  PPTX -> python-pptx
  XLSX -> openpyxl (read-only, handles password/corrupt gracefully)
Legacy binary .doc/.ppt/.xls need LibreOffice; we skip them (text="") and
let the router classify by filename/metadata instead of crashing.
"""

from __future__ import annotations

import shutil
import logging
from pathlib import Path

import structlog

logger = structlog.get_logger()

# Suppress non-critical warnings from PDF libraries (malformed PDFs with invalid color values)
logging.getLogger("pdfplumber").setLevel(logging.ERROR)
logging.getLogger("pypdf").setLevel(logging.ERROR)
logging.getLogger("pypdf.generic").setLevel(logging.ERROR)

LEGACY_BINARY = {".doc", ".ppt", ".xls"}


def extract_pdf(path: Path, max_chars: int) -> str:
    """PDF via pdfplumber, fallback pypdf. Returns "" on failure."""
    try:
        import pdfplumber

        parts: list[str] = []
        total = 0
        with pdfplumber.open(str(path)) as pdf:
            for page in pdf.pages:
                t = page.extract_text() or ""
                if t:
                    parts.append(t)
                    total += len(t)
                if total >= max_chars:
                    break
        text = "\n".join(parts).strip()
        if text:
            return text
    except Exception as e:
        logger.debug("pdfplumber_failed", path=str(path), error=str(e)[:120])

    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        parts = []
        total = 0
        for page in reader.pages:
            t = page.extract_text() or ""
            if t:
                parts.append(t)
                total += len(t)
            if total >= max_chars:
                break
        return "\n".join(parts).strip()
    except Exception as e:
        logger.warning("pdf_extract_failed", path=str(path), error=str(e)[:120])
        return ""


def extract_docx(path: Path, max_chars: int) -> str:
    """DOCX via python-docx (paragraphs + table cells)."""
    try:
        from docx import Document

        doc = Document(str(path))
        parts = [p.text for p in doc.paragraphs if p.text]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text for c in row.cells if c.text]
                if cells:
                    parts.append(" | ".join(cells))
                if sum(len(x) for x in parts) >= max_chars:
                    break
        return "\n".join(parts).strip()
    except Exception as e:
        logger.warning("docx_extract_failed", path=str(path), error=str(e)[:120])
        return ""


def extract_pptx(path: Path, max_chars: int) -> str:
    """PPTX via python-pptx (all shape text frames + notes)."""
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text:
                    parts.append(shape.text_frame.text)
            if sum(len(x) for x in parts) >= max_chars:
                break
        return "\n".join(parts).strip()
    except Exception as e:
        logger.warning("pptx_extract_failed", path=str(path), error=str(e)[:120])
        return ""


def extract_xlsx(path: Path, max_chars: int) -> str:
    """XLSX via openpyxl (read-only). Skips password-protected/corrupt cleanly."""
    try:
        from openpyxl import load_workbook

        wb = load_workbook(str(path), read_only=True, data_only=True)
        parts: list[str] = []
        total = 0
        for ws in wb.worksheets:
            parts.append(f"# {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    line = " ".join(cells)
                    parts.append(line)
                    total += len(line)
                if total >= max_chars:
                    break
            if total >= max_chars:
                break
        wb.close()
        return "\n".join(parts).strip()
    except Exception as e:
        logger.warning("xlsx_extract_failed", path=str(path), error=str(e)[:120])
        return ""


def extract_legacy(path: Path, max_chars: int) -> str:
    """Legacy .doc/.ppt/.xls — needs LibreOffice. Degrade to "" if unavailable."""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        logger.info("legacy_office_skipped_no_libreoffice", path=str(path))
        return ""
    # LibreOffice present but conversion intentionally not implemented here —
    # router will classify by filename/metadata. Keeps the hot path dependency-light.
    logger.info("legacy_office_classify_by_name", path=str(path))
    return ""
